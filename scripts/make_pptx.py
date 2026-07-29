"""Build ``deck.pptx``: the talk as an editable PowerPoint.

The HTML deck is what we present from, but a submission wants a file anyone can
open and edit, so this writes the same seven slides as real PowerPoint text boxes
and pictures rather than screenshots. Every string here is text you can retype in
PowerPoint, and every image comes from the repository's own output.

    python scripts/make_pptx.py            # -> deck.pptx

Type is deliberately large: 16:9 at 13.33 by 7.5 inches, headlines at 30 pt and
body copy at 17 pt, because a judge reads this from across a room rather than on a
laptop. If a slide overflows, cut words rather than shrinking the type.
"""

from __future__ import annotations

import argparse
import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

OUT = "deck.pptx"
VAL_RUN = "results/val/val_w40_sl0.0.json"

# The deck's own palette, so the PowerPoint and the HTML agree.
GREEN = RGBColor(0x0B, 0x53, 0x40)      # brand, chrome only
LIME = RGBColor(0xD6, 0xEF, 0x4A)
INK = RGBColor(0x14, 0x20, 0x1B)
MID = RGBColor(0x3D, 0x4A, 0x44)
SOFT = RGBColor(0x6B, 0x77, 0x72)
PAPER = RGBColor(0xEF, 0xEF, 0xE9)
TRUE = RGBColor(0x2A, 0x78, 0xD6)       # ground truth
BASE = RGBColor(0xEB, 0x68, 0x34)       # distortion-optimal
OURS = RGBColor(0x1B, 0xAF, 0x7A)       # tumor-aware

W, H = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.72)

SANS = "Helvetica Neue"
MONO = "Menlo"


def _bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _text(slide, x, y, w, h, runs, size=17, color=MID, font=SANS, bold=False,
          align=PP_ALIGN.LEFT, spacing=1.15, caps=False, tracking=False):
    """One text box. ``runs`` is a string or a list of (text, bold, colour)."""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = spacing
    items = [(runs, bold, color)] if isinstance(runs, str) else runs
    for text, b, c in items:
        r = p.add_run()
        r.text = text.upper() if caps else text
        r.font.size = Pt(size)
        r.font.bold = b
        r.font.name = font
        r.font.color.rgb = c
    return box


def _rule(slide, y, color=RGBColor(0xDB, 0xDB, 0xD2)):
    line = slide.shapes.add_shape(1, MARGIN, y, W - 2 * MARGIN, Emu(9525))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    line.shadow.inherit = False
    return line


def _header(slide, section):
    """The running head: section on the left, deck name on the right."""
    _text(slide, MARGIN, Inches(0.42), Inches(7), Inches(0.4), section, size=12,
          color=GREEN, font=MONO, caps=True)
    _text(slide, W - MARGIN - Inches(4), Inches(0.42), Inches(4), Inches(0.4),
          "Tumor-aware SR  ·  ACVSS", size=12, color=SOFT, font=MONO,
          align=PP_ALIGN.RIGHT)
    _rule(slide, Inches(0.92))


def _headline(slide, text, y=Inches(1.18), size=30):
    return _text(slide, MARGIN, y, W - 2 * MARGIN - Inches(0.5), Inches(1.5), text,
                 size=size, color=INK, bold=True, spacing=1.06)


def _card(slide, x, y, w, h, title, body, rail=OURS):
    """A definition card: a colour rail, a term, and what it means."""
    bar = slide.shapes.add_shape(1, x, y, Inches(0.055), h)
    bar.fill.solid(); bar.fill.fore_color.rgb = rail
    bar.line.fill.background(); bar.shadow.inherit = False
    _text(slide, x + Inches(0.28), y, w - Inches(0.36), Inches(0.5), title,
          size=20, color=INK, bold=True)
    _text(slide, x + Inches(0.28), y + Inches(0.55), w - Inches(0.36),
          h - Inches(0.6), body, size=16, color=MID, spacing=1.25)


def _stat(slide, x, y, w, number, label, note, color=INK, size=54):
    _text(slide, x, y, w, Inches(0.95), number, size=size, color=color, bold=True,
          font=MONO)
    _text(slide, x, y + Inches(0.95), w, Inches(0.4), label, size=17, color=INK,
          bold=True)
    _text(slide, x, y + Inches(1.35), w, Inches(1.0), note, size=15, color=SOFT,
          spacing=1.2)


def _picture(slide, path, x, y, w=None, h=None):
    if not os.path.exists(path):
        return None
    return slide.shapes.add_picture(path, x, y, width=w, height=h)


# --------------------------------------------------------------------------- #
# the seven slides                                                            #
# --------------------------------------------------------------------------- #

NAMES = ("Adiza Alhassan  ·  Nthabiseng Thema  ·  Albert Dodoo  ·  "
         "Victor Oyindouye Miene  ·  Hassan Suliman")


def title(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, GREEN)
    _text(s, MARGIN, Inches(2.05), Inches(9), Inches(0.4),
          "ACVSS  ·  MRI super-resolution safety", size=14, color=LIME, font=MONO,
          caps=True)
    _text(s, MARGIN, Inches(2.55), Inches(11), Inches(1.5),
          [("When sharper means ", True, RGBColor(0xFF, 0xFF, 0xFF)),
           ("blind", True, LIME), (".", True, RGBColor(0xFF, 0xFF, 0xFF))],
          size=54, spacing=1.05)
    _text(s, MARGIN, Inches(4.05), Inches(9.4), Inches(1.4),
          [("Super-resolution makes a cheap low-field brain scan look crisp. "
            "Trained only for image quality, it can quietly ", False,
            RGBColor(0xD8, 0xE4, 0xDF)),
           ("erase a small tumor", True, RGBColor(0xFF, 0xFF, 0xFF)),
           (". We measure how often that happens, and fix it with a tumor-aware "
            "objective.", False, RGBColor(0xD8, 0xE4, 0xDF))],
          size=19, spacing=1.35)
    _text(s, MARGIN, Inches(5.75), Inches(11.5), Inches(0.8), NAMES, size=16,
          color=RGBColor(0xB9, 0xCE, 0xC6), font=MONO, spacing=1.4)
    return s


def access(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, PAPER)
    _header(s, "01  Why this matters")
    _headline(s, "Cheap, low-quality scanners are the only realistic way to widen "
                 "MRI access.")
    _picture(s, "figures/mri_room.jpg", MARGIN, Inches(2.6), w=Inches(6.6))
    _text(s, MARGIN, Inches(7.0), Inches(6.6), Inches(0.5),
          "A three tesla scanner needs a shielded room, a cooling supply and power "
          "that does not fail.", size=13, color=SOFT)
    x = MARGIN + Inches(7.2)
    w = W - x - MARGIN
    _stat(s, x, Inches(2.55), w, "14", "Ghana.",
          "For more than 30 million people, and two thirds of them sit in Greater "
          "Accra.", color=BASE)
    _stat(s, x, Inches(4.15), w, "<1", "Per million, much of sub-Saharan Africa.",
          "MRI is scarce where the disease burden is not.", color=BASE)
    _stat(s, x, Inches(5.75), w, "37", "Per million, high-income countries.",
          "Up to this many, for the same imaging need.")
    return s


def contribution(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, PAPER)
    _header(s, "02  Contribution")
    _headline(s, "Penalise the model heavily for mistakes where the tumor is.")
    y, h = Inches(2.9), Inches(2.5)
    w = (W - 2 * MARGIN - Inches(0.8)) / 2
    _card(s, MARGIN, y, w, h, "How we degrade",
          "We take a real high-field scan and throw away the fine detail a cheap "
          "scanner cannot capture, then add noise. Every blurry input still has its "
          "true scan to be judged against.", rail=TRUE)
    _card(s, MARGIN + w + Inches(0.8), y, w, h, "The tumor-aware loss",
          "A mistake inside the tumor costs the model far more than the same mistake "
          "in healthy tissue. It can no longer buy score by smoothing a lesion away.",
          rail=OURS)
    _text(s, MARGIN, Inches(5.75), W - 2 * MARGIN, Inches(0.6),
          "Same network, same data. The only thing we changed is what it is punished "
          "for.", size=16, color=SOFT)
    return s


def method(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, PAPER)
    _header(s, "03  Method")
    _headline(s, "Degrade a real scan, reconstruct it two ways, and ask one frozen "
                 "segmentation network what it can still find.", size=27)
    _picture(s, "figures/architecture.png", MARGIN, Inches(3.1),
             w=W - 2 * MARGIN)
    _text(s, MARGIN, Inches(6.6), W - 2 * MARGIN, Inches(0.6),
          "No GAN, deliberately: an adversarial loss rewards inventing plausible "
          "tissue. 17,233 slices, 468 patients.", size=16, color=SOFT)
    return s


def result(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, PAPER)
    _header(s, "04  Result  ·  real BraTS, 70 unseen patients")
    _headline(s, "Super-resolution recovers tumor the degradation destroys. "
                 "Ours recovers more of it.")
    x = MARGIN
    w = (W - 2 * MARGIN - Inches(1.0)) / 3
    for num, label, note, col in [
            ("62.2%", "The degraded scan", "missed before any reconstruction", SOFT),
            ("58.0%", "Standard SR (baseline)", "recovers some of it", BASE),
            ("51.3%", "Tumor-aware SR (ours)", "6.7 points better, at matched quality",
             OURS)]:
        _stat(s, x, Inches(2.65), w, num, label, note, color=col, size=44)
        x += w + Inches(0.5)
    _picture(s, "figures/floor_by_size.png", MARGIN, Inches(4.5), h=Inches(2.5))
    _text(s, MARGIN + Inches(6.4), Inches(4.7), W - MARGIN - Inches(6.9),
          Inches(2.2),
          "Grey is what the frozen segmenter misses on the untouched scan. The block "
          "above it is what each objective adds: standard SR in orange, ours in "
          "green.", size=16, color=MID, spacing=1.25)
    return s


def quality(prs):
    """The table that earns the phrase "at matched quality"."""
    import json
    import os

    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, PAPER)
    _header(s, "04  Result  ·  image quality")
    _headline(s, "The two reconstructions are level on every quality metric we are "
                 "judged by.")

    rows = [("low-res input", "lowres", SOFT),
            ("standard SR (baseline)", "distortion", BASE),
            ("tumor-aware SR (ours)", "tumor-aware", OURS)]
    res = {}
    if os.path.exists(VAL_RUN):
        res = json.load(open(VAL_RUN))["results"]

    cols = [Inches(4.6), Inches(2.4), Inches(2.0), Inches(2.4)]
    heads = ["image", "PSNR (dB), in brain", "SSIM", "Dice, segmenter vs truth"]
    y = Inches(2.9)
    x = MARGIN
    for w, head in zip(cols, heads):
        _text(s, x, y, w, Inches(0.5), head, size=14, color=SOFT, font=MONO,
              caps=True)
        x += w
    _rule(s, y + Inches(0.52))

    y += Inches(0.72)
    for label, key, hue in rows:
        r = res.get(key)
        if not r:
            continue
        vals = [label, f"{r['psnr_brain']:.2f}", f"{r['ssim']:.3f}",
                f"{r['dice']:.3f}"]
        x = MARGIN
        for w, v, first in zip(cols, vals, [True, False, False, False]):
            _text(s, x, y, w, Inches(0.55), v, size=19,
                  color=hue if first else INK, bold=first,
                  font=SANS if first else MONO)
            x += w
        _rule(s, y + Inches(0.6))
        y += Inches(0.82)

    _text(s, MARGIN, y + Inches(0.3), W - 2 * MARGIN, Inches(0.8),
          "70 unseen patients. If these rows were not level, the erasure difference "
          "would be a quality difference wearing a safety costume.", size=16,
          color=SOFT, spacing=1.25)
    return s


def viewports(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, PAPER)
    _header(s, "05  Result  ·  four viewports")
    _headline(s, "The lesions the baseline loses show up in 3D as empty blue shells.")
    panels = [("brain3d_true.png", "Ground truth", "the true lesions, in blue"),
              ("brain3d_tumor_aware.png", "Tumor-aware (ours)", "lesions preserved"),
              ("brain3d_distortion.png", "Baseline", "small lesions dropped"),
              ("brain3d_uncertainty.png", "Uncertainty", "where the model is unsure")]
    gap = Inches(0.28)
    w = (W - 2 * MARGIN - 3 * gap) / 4
    x = MARGIN
    for path, name, note in panels:
        _picture(s, path, x, Inches(2.85), w=w)
        _text(s, x, Inches(2.85) + w + Inches(0.12), w, Inches(0.4), name, size=17,
              color=INK, bold=True)
        _text(s, x, Inches(2.85) + w + Inches(0.5), w, Inches(0.6), note, size=14,
              color=SOFT)
        x += w + gap
    _text(s, MARGIN, Inches(6.75), W - 2 * MARGIN, Inches(0.5),
          "A blue shell with nothing inside it is a lesion the model lost.", size=16,
          color=SOFT)
    return s


def next_steps(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, PAPER)
    _header(s, "06  Next steps")
    _headline(s, "Next: generative models, tested with the readout we just built.")
    y, h = Inches(2.9), Inches(2.4)
    w = (W - 2 * MARGIN - Inches(1.0)) / 3
    x = MARGIN
    for title_, body, rail in [
            ("Why generative",
             "Diffusion and adversarial models make the sharpest low-field images "
             "anyone has produced. They are also the ones most likely to invent "
             "tissue.", BASE),
            ("Why we can now try them",
             "We can measure what they lose and what they fabricate. Sharpness no "
             "longer has to be taken on trust.", OURS),
            ("Then the one number",
             "A single evaluation on the 94 patients nothing has touched.", SOFT)]:
        _card(s, x, y, w, h, title_, body, rail=rail)
        x += w + Inches(0.5)
    _text(s, MARGIN, Inches(5.7), W - 2 * MARGIN, Inches(0.6),
          "A safety metric is what makes a generative model testable rather than "
          "impressive.", size=16, color=SOFT)
    return s


def build(out: str = OUT) -> str:
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    for fn in (title, access, contribution, next_steps, method, result, quality,
               viewports):
        fn(prs)
    prs.save(out)
    print(f"wrote {out} ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
    return out


def main() -> None:
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("--out", default=OUT)
    build(ap.parse_args().out)


if __name__ == "__main__":
    main()
